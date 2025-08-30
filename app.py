import os
from dotenv import load_dotenv
import streamlit as st
from PIL import Image
import google.generativeai as genai
import docx
import pptx
import openpyxl
import io

# --------------------🔐 Load and Configure API --------------------
load_dotenv()
try:
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        st.error("GOOGLE_API_KEY not found. Please set it in your .env file.")
        st.stop()
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("gemini-1.5-flash")
except Exception as e:
    st.error(f"Error configuring API: {e}")
    st.stop()

# --------------------🛠️ Helper Functions for File Processing --------------------
def process_uploaded_file(uploaded_file):
    """Detects file type and extracts content (text or image)."""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    
    try:
        if file_extension in ['.jpg', '.jpeg', '.png']:
            return Image.open(uploaded_file)
        elif file_extension == '.pdf':
            from pypdf import PdfReader
            text = ""
            reader = PdfReader(io.BytesIO(uploaded_file.getvalue()))
            for page in reader.pages:
                text += page.extract_text() or ""
            return text
        elif file_extension == '.docx':
            doc = docx.Document(io.BytesIO(uploaded_file.getvalue()))
            return "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == '.pptx':
            prs = pptx.Presentation(io.BytesIO(uploaded_file.getvalue()))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif file_extension == '.xlsx':
            workbook = openpyxl.load_workbook(io.BytesIO(uploaded_file.getvalue()))
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"--- Sheet: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    text += " | ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            return text
    except Exception as e:
        st.error(f"Error processing file: {e}")
        return None
    return None

# --------------------🎨 App UI and Layout --------------------
st.set_page_config(page_title="Gemini AI Chatbot", page_icon="🤖", layout="wide")

st.markdown(
    "<h2 style='text-align: center;'>🤖 Gemini AI Chatbot 💬</h2>",
    unsafe_allow_html=True
)

# --- Sidebar for file uploads ---
with st.sidebar:
    st.header("Upload for Context")
    uploaded_file = st.file_uploader(
        "Upload an image, PDF, DOCX, XLSX, or PPTX file",
        type=['png', 'jpg', 'jpeg', 'pdf', 'docx', 'xlsx', 'pptx']
    )
    
    if uploaded_file:
        with st.spinner("Processing file..."):
            file_content = process_uploaded_file(uploaded_file)
            st.session_state.file_context = file_content
            if isinstance(file_content, Image.Image):
                st.image(file_content, caption="Uploaded Image")
            st.success("File processed! Ask a question about it.")

# --------------------🧠 Memory and Chat History --------------------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "file_context" not in st.session_state:
    st.session_state.file_context = None

# --- Display previous chat messages ---
for message in st.session_state.chat_history:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# --------------------💬 Chat Input and Response Logic --------------------
if user_input := st.chat_input("💡 Ask me anything..."):
    # Display user message
    st.session_state.chat_history.append({"role": "user", "content": user_input})
    with st.chat_message("user"):
        st.markdown(user_input)
    
    # --- Generate AI response ---
    with st.spinner("🧠 Thinking..."):
        
        # ** DYNAMIC PROMPT LOGIC **
        if st.session_state.file_context is not None:
            # If there's a file, use the strict, context-aware prompt
            system_prompt = "You are an expert assistant. Your task is to answer the user's question based ONLY on the provided file context. Do not use any external knowledge. If the answer is not in the context, you must state: 'I don't know, as the answer is not in the provided document.'"
            model_input = [system_prompt, "CONTEXT FROM UPLOADED FILE:", st.session_state.file_context, "USER'S QUESTION:", user_input]
            # Clear file context after using it once
            st.session_state.file_context = None 
        else:
            # If no file, use a general conversational prompt
            system_prompt = "You are a helpful and friendly AI assistant. Answer the user's questions. Be friendly and conversational."
            model_input = [system_prompt, user_input]

        try:
            # The chat history is passed to the model for context
            chat = model.start_chat(history=[
                {"role": "user", "parts": [msg["content"]]} if msg["role"] == "user" else {"role": "model", "parts": [msg["content"]]}
                for msg in st.session_state.chat_history[:-1] # Exclude the current user input
            ])
            response = chat.send_message(model_input)
            bot_reply = response.text

            st.chat_message("assistant").markdown(bot_reply)
            st.session_state.chat_history.append({"role": "assistant", "content": bot_reply})

        except Exception as e:
            bot_reply = f"Sorry, an error occurred: {e}"
            st.error(bot_reply)
            st.session_state.chat_history.append({"role": "assistant", "content": bot_reply})


# --- Creator Attribution ---
st.markdown(
    """
    <style>
        .footer {
            position: fixed;
            right: 10px;
            bottom: 10px;
            width: auto;
            background-color: transparent;
            color: #aaa;
            text-align: right;
            padding: 5px;
            font-size: 12px;
        }
    </style>
    <div class="footer">
        Created BY- Ridham Vashishth
    </div>
    """,
    unsafe_allow_html=True
)

